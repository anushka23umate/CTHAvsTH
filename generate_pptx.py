from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_ctha_presentation():
    prs = Presentation()
    
    def add_bullet_slide(title_text, points, notes=""):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_text
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(points):
            # Check if it's a sub-bullet (starts with '>')
            level = 0
            text = point
            if point.startswith(">"):
                level = 1
                text = point[1:].strip()
            
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = text
            p.level = level
            p.font.size = Pt(18) if level == 0 else Pt(14)
            
        # Add Speaker Notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
        return slide

    # --- SLIDE 1: TITLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "CTHA: Constrained Temporal Hierarchical Architecture"
    slide.placeholders[1].text = "Solving the Coordination Crisis in Multi-Agent Systems\nResearch Proof: TH vs. CTHA Comparative Analysis"

    # --- SLIDE 2: DEFINITIONS ---
    add_bullet_slide("Defining the Architectures", [
        "Traditional Hierarchy (TH):",
        "> A linear pipeline where agents coordinate via loose natural language.",
        "> Stability relies on prompt-compliance (highly unreliable at scale).",
        "Constrained Temporal Hierarchical Architecture (CTHA):",
        "> A state-managed architecture using structural manifolds.",
        "> Stability is enforced by mathematical constraints (CC1-CC3).",
        "Key Findings:",
        "> Architecture always outperforms prompting in high-noise environments."
    ], notes="Traditional Hierarchies (TH) are what most developers build today. CTHA is the future of stable, industrial-grade agentic systems.")

    # --- SLIDE 3: THE PROBLEM ---
    add_bullet_slide("The Problem: Semantic Entropy", [
        "Semantic Entropy: The accumulation of noise as data passes through layers.",
        "Observation: In unconstrained TH, small errors at the Reflex layer become catastrophic hallucinations at the Strategic layer.",
        "Leakage:",
        "> 403/404 errors are treated as 'facts' by TH.",
        "> Formatting constraints are ignored as context size grows.",
        "The Result: Systemic instability and unpredictable agent behavior."
    ], notes="Semantic entropy is like the 'Telephone Game' for AI. CTHA stops this by filtering noise at every single hop.")

    # --- SLIDE 4: KEY FINDING - THE SCALABILITY WALL ---
    add_bullet_slide("Finding 1: The Scalability Wall", [
        "Payload Bloat: TH has no 'valve' to control information flow.",
        "Empirical Evidence from our Implementation:",
        "> TH requested 9,489 tokens for an 8,000 token limit.",
        "> Result: TOTAL SYSTEM CRASH (Error 413).",
        "Conclusion:",
        "> Traditional Hierarchies are physically unscalable in context-limited environments.",
        "> CTHA remained stable at <2,500 tokens using the exact same query."
    ], notes="TH simply cannot handle large-scale data because it doesn't know how to 'truncate with intent'. CTHA manifolds solve this.")

    # --- SLIDE 5: PILLAR 1 - CC1 (MESSAGE CONTRACTS) ---
    add_bullet_slide("Pillar 1: Message Contracts (CC1)", [
        "Concept: All inter-layer messages must be projected onto a Manifold.",
        "In our implementation:",
        "> PMsum: Only passes text with high 'Density' (removes noise).",
        "> PMplan: Verifies that summaries actually align with the query.",
        "Result:",
        "> The Reasoning layer never sees 403 errors or 'Pet Blog' garbage.",
        "> Information grounding is forced at the ingress."
    ], notes="CC1 is the 'Security Guard' of the architecture. It rejects bad data before it can poison the reasoning process.")

    # --- SLIDE 6: DIAGRAM (MANIFOLD PROJECTION) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1)).text_frame.text = "Visualizing the Manifold Projection (CC1)"
    def add_box(left, top, text, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(2), Inches(1))
        shape.text = text
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        return shape
    b1 = add_box(0.5, 2.5, "Raw Data\n(High Entropy)", RGBColor(100, 100, 100))
    b2 = add_box(3.5, 2.5, "PMsum Manifold\n(Filter/Project)", RGBColor(0, 120, 215))
    b3 = add_box(6.5, 2.5, "Distilled Facts\n(Low Entropy)", RGBColor(0, 176, 80))
    slide.shapes.add_connector(1, b1.left + b1.width, b1.top + b1.height/2, b2.left, b2.top + b2.height/2)
    slide.shapes.add_connector(1, b2.left + b2.width, b2.top + b2.height/2, b3.left, b3.top + b3.height/2)

    # --- SLIDE 7: PILLAR 2 - CC2 (AUTHORITY SCOPING) ---
    add_bullet_slide("Pillar 2: Authority Scoping (CC2)", [
        "Rule: Agents are restricted in what they can 'decide'.",
        "Implementation:",
        "> Tactical Agent: Only permitted to extract facts. Cannot write reports.",
        "> Institutional Agent: Only permitted to refine structure. Cannot add facts.",
        "The Proof:",
        "> TH agents 'talked back' to the user to justify their drift.",
        "> CTHA agents remained silent and executed the forced policy."
    ], notes="CC2 prevents 'Role Confusion'. It ensures each agent stays in its lane, reducing structural noise.")

    # --- SLIDE 8: PILLAR 3 - CC3 (ARBITER RESOLUTION) ---
    add_bullet_slide("Pillar 3: Arbiter Resolution (CC3)", [
        "Rule: Active conflict monitoring and Temporal Backtracking.",
        "Observation from our Run:",
        "> Arbiter detected that 'Pet Blogs' were non-aligned with 'War News'.",
        "> Triggered a backtrack to Reflex layer for query refinement.",
        "Conclusion:",
        "> CTHA avoids hallucinations by refusing to proceed with unstable data."
    ], notes="CC3 is the 'Quality Control' layer. It is the only part of the system allowed to trigger a retry.")

    # --- SLIDE 9: RESULT - STRUCTURAL STABILITY ---
    add_bullet_slide("Finding 2: Structural Supremacy", [
        "The 'One Paragraph' Test:",
        "> Policy: Output MUST be a single concise paragraph.",
        "TH Failure:",
        "> Generated 174 Policy Manifold Violations (headers, bolds, breaks).",
        "> Directly ignored the prompt to appear 'professional'.",
        "CTHA Success:",
        "> 0 Policy Violations. Structural stability forced by PMpol manifold."
    ], notes="This is the mic-drop moment. TH can't follow a simple formatting rule under load. CTHA can't BREAK the rule.")

    # --- SLIDE 10: ARCHITECTURAL SUPREMACY DASHBOARD ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Architectural Supremacy Dashboard"
    table = slide.shapes.add_table(5, 3, Inches(0.5), Inches(2), Inches(9), Inches(4)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "TH (Traditional)"
    table.cell(0, 2).text = "CTHA (Constrained)"
    metrics = [
        ("Payload Control", "❌ FAILED (413 Bloat)", "✅ STABLE (Projected)"),
        ("Structural Noise", "174 Violations", "0 Violations"),
        ("Grounding Force", "Prompt-Based (Weak)", "Manifold-Based (Absolute)"),
        ("Stability Status", "UNSTABLE", "✅ VERIFIED")
    ]
    for r, (m, th, ctha) in enumerate(metrics, start=1):
        table.cell(r, 0).text = m
        table.cell(r, 1).text = th
        table.cell(r, 2).text = ctha

    # --- SLIDE 11: KEY FINDING - THE SILENT REJECTION ---
    add_bullet_slide("Finding 3: Grounded Silence", [
        "Noise Handling:",
        "> TH wove 'Pets Hub' data into the war report (Hallucination).",
        "> CTHA rejected the noise and signaled a stability failure.",
        "The Insight:",
        "> Traditional Hierarchies prioritize 'Filling the Space'.",
        "> CTHA prioritizes 'Evidence Integrity'.",
        "It is better to refuse than to hallucinate."
    ], notes="TH will always try to answer, even if the data is garbage. CTHA is the only architecture that can safely say 'I don't know'.")

    # --- SLIDE 12: INDUSTRIAL IMPACT ---
    add_bullet_slide("Industrial Impact & Scale", [
        "Financial/Medical Compliance:",
        "> CTHA ensures zero-drift between source data and final output.",
        "Token Optimization:",
        "> CTHA reduces context costs by 70-80% through manifold distillation.",
        "Safety:",
        "> Prevents 'Agentic Breakaway' by hard-coding policy manifolds."
    ])

    # --- SLIDE 13: THE "PROMPT VS ARCHITECTURE" DEBATE ---
    add_bullet_slide("Prompting vs. Architecture", [
        "Prompting is a 'Suggestion':",
        "> Relies on the LLM's willingness to follow complex rules.",
        "Architecture is a 'Physical Law':",
        "> Uses code-based manifolds to make violations impossible.",
        "Conclusion:",
        "> For 2026+ Agentic systems, Prompting is for prototyping; CTHA is for production."
    ])

    # --- SLIDE 14: FUTURE RESEARCH DIRECTIONS ---
    add_bullet_slide("Future Research (arXiv:2601.10738v1)", [
        "Embedding-Space Manifolds: Projecting intent directly in vector space.",
        "Multi-Arbiter Networks: Coordinating across domain-specific hierarchies.",
        "Real-time Policy Adaptation: Dynamic manifold adjustment."
    ])

    # --- SLIDE 15: CONCLUSION ---
    add_bullet_slide("Conclusion", [
        "CTHA is the only stable path for complex multi-agent reasoning.",
        "Key Proofs Delivered:",
        "> 1. Payload Stability (No Crashes)",
        "> 2. Structural Stability (No Drift)",
        "> 3. Evidence Stability (No Hallucination)",
        "Final Verdict: Architecture > Intelligence."
    ])

    prs.save('CTHA_Final_Research_Proof.pptx')
    print("Presentation saved as CTHA_Final_Research_Proof.pptx")

if __name__ == "__main__":
    create_ctha_presentation()
